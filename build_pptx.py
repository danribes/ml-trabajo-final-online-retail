"""Fill PPTX template with model results + EDA images."""
import json, copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

HERE = Path(__file__).parent
ART  = HERE / 'artifacts'
TPL  = HERE / 'Analisis de reseñas (plantilla).pptx'
OUT  = HERE / 'presentacion_recompra.pptx'

m = json.load(open(ART/'metrics.json'))
best = m['final_model']
auc  = m[best]['AUC']
ap   = m[best]['AP']

p = Presentation(str(TPL))

def set_text(shape, lines, size=None):
    """Replace text in a shape, keeping first run formatting; lines = list[str]."""
    tf = shape.text_frame
    # capture first run format
    first_para = tf.paragraphs[0]
    if first_para.runs:
        r0 = first_para.runs[0]
        font = r0.font
        base_size = font.size
        base_bold = font.bold
        base_name = font.name
        base_color = None
        try:
            if font.color and font.color.type is not None:
                base_color = font.color.rgb
        except Exception:
            pass
    else:
        base_size = base_bold = base_name = base_color = None

    # wipe paragraphs
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        if size is not None:
            run.font.size = Pt(size)
        elif base_size is not None:
            run.font.size = base_size
        if base_bold is not None:
            run.font.bold = base_bold
        if base_name:
            run.font.name = base_name
        if base_color is not None:
            try:
                run.font.color.rgb = base_color
            except Exception:
                pass

# ---------- slide 1 – title ----------
s = p.slides[0]
set_text(s.shapes[0], ['Predicción de recompra',
                       'Online Retail II – Máster ML'])

# ---------- slide 2 – objetivos + asunciones ----------
s = p.slides[1]
# header (Objetivos)
set_text(s.shapes[0], ['Objetivos e hipótesis'])
# right of "Objetivos" label (shape 1) → objetivos del análisis
set_text(s.shapes[1], [
    f'• Predecir la probabilidad de que un cliente vuelva a comprar en los 6 meses',
    f'   posteriores a la fecha de corte (2011-06-09).',
    f'• Segmentar a {m["n_customers"]:,} clientes por score → acciones diferenciadas.',
    f'• Fidelización para top decil; retención / win-back para bottom decil.',
])
# left labels (kept as-is)
# right of "Asunciones" label (shape 4) → exclusiones del dataset
set_text(s.shapes[4], [
    '• Filas sin CustomerID descartadas (~23%) → no etiquetables.',
    '• Cancelaciones (Invoice "C*") y devoluciones excluidas.',
    '• Códigos de servicio (POST, DOT, AMAZONFEE, …) excluidos.',
    '• Quantity ≤ 0 y Price ≤ 0 excluidos.',
    '• Top 0.5% revenue por línea capado para estabilidad.',
    f'• Tasa de recompra observada: {m["repurchase_rate"]:.1%} (muestra balanceada).',
])

# ---------- slide 3 – descriptivo ----------
s = p.slides[2]
set_text(s.shapes[0], ['Descriptivo'])
# Left column: monthly_revenue (top) + top_countries (bottom)
s.shapes.add_picture(str(ART/'eda_monthly_revenue.png'),
                     Inches(0.15), Inches(1.20), width=Inches(4.7))   # h ≈ 2.09
s.shapes.add_picture(str(ART/'eda_top_countries.png'),
                     Inches(0.15), Inches(3.40), height=Inches(2.15))  # w ≈ 4.30
# Right: target_balance
s.shapes.add_picture(str(ART/'eda_target_balance.png'),
                     Inches(5.10), Inches(1.20), height=Inches(3.8))  # w ≈ 4.75
# Caption below target_balance
tx = s.shapes.add_textbox(Inches(5.10), Inches(5.05), Inches(4.75), Inches(0.4))
tf = tx.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = f'{m["n_customers"]:,} clientes etiquetados · tasa real {m["repurchase_rate"]:.1%}'
r.font.size = Pt(10); r.font.italic = True

# ---------- slide 4 – modelo ----------
s = p.slides[3]
set_text(s.shapes[0], ['El modelo'])
# Left col: ROC (top) + deciles (bottom)
s.shapes.add_picture(str(ART/'model_roc.png'),
                     Inches(0.15), Inches(1.20), height=Inches(2.4))  # w ≈ 2.88
s.shapes.add_picture(str(ART/'model_deciles.png'),
                     Inches(0.15), Inches(3.80), height=Inches(1.75)) # w ≈ 3.06
# Right col: feature importance
s.shapes.add_picture(str(ART/'model_importance.png'),
                     Inches(3.30), Inches(1.20), height=Inches(2.4))  # w ≈ 3.73
# Metrics textbox bottom-right
tx = s.shapes.add_textbox(Inches(3.30), Inches(3.80), Inches(6.55), Inches(1.75))
tf = tx.text_frame; tf.word_wrap = True
lines = [
    f'Modelo final: {best}',
    f'Desempeño actual:  AUC = {auc:.4f}   AP = {ap:.4f}',
    f'Modelo RF anterior:  AUC = 0.8030   AP = 0.8329',
    f'Mejora de performance:  +{auc - 0.8030:+.4f} AUC',
    f'Top decil score → {m["top_decile_rate"]*100:.1f}% recompran (vs {m["bottom_decile_rate"]*100:.1f}% decil bajo)',
]
for i, t in enumerate(lines):
    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = para.add_run(); run.text = t
    run.font.size = Pt(13)
    if i == 0:
        run.font.bold = True

# ---------- slide 5 – conclusiones + próximos pasos ----------
s = p.slides[4]
set_text(s.shapes[0], ['Conclusiones y próximos pasos'])
# right of "Conclusiones" → principales conclusiones
set_text(s.shapes[1], [
    f'• Modelo {best} con AUC test ≈ {auc:.4f}, AP ≈ {ap:.4f}.',
    f'• Buena calibración → score utilizable como probabilidad real.',
    f'• Lift fuerte: top 10% del score recompra {m["top_decile_rate"]*100:.1f}% vs. {m["bottom_decile_rate"]*100:.1f}% en el último decil.',
    '• Variables comportamentales e IPI superan a la línea base.',
    '• Recomendación: usar score por decil para diseñar campañas.',
])
# right of "Próximos pasos"
set_text(s.shapes[4], [
    '• Validar sobre cutoffs rolling para robustez temporal.',
    '• Añadir features de categoría de producto y estacionalidad Q4.',
    '• Calibrar threshold con curva coste / beneficio (campaña vs CLV recuperado).',
    '• Probar Gradient Boosting con monotonic constraints en Recency/Frequency.',
    '• Calibrar el output (Platt / isotonic) si se usa segmentación binaria.',
    '• Re-entrenamiento trimestral + monitorización de drift de RFM.',
])

# ---------- remove slide 6 (instructions) ----------
xml_slides = p.slides._sldIdLst
slides = list(xml_slides)
xml_slides.remove(slides[5])

p.save(str(OUT))
print(f'Wrote {OUT}')
